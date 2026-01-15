from __future__ import annotations

import asyncio
import tempfile
from pathlib import Path

import pytest

from src.processors.office_processor import OfficeProcessor


def _make_pptx(tmp: Path) -> Path:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        pytest.skip(f"python-pptx nicht verfügbar: {e}")

    p = tmp / "sample.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])  # title + content
    slide.shapes.title.text = "Titel"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Hallo"
    pres.save(str(p))
    return p


def test_office_processor_pptx_creates_markdown() -> None:
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        pptx_path = _make_pptx(tmp)

        proc = OfficeProcessor(process_id="test-pptx")
        res = asyncio.run(
            proc.process(
                pptx_path,
                include_images=True,
                include_previews=True,
                use_cache=False,
                base_cache_dir=tmp / "cache",
            )
        )

        assert res.data.metadata.format == "pptx"
        assert "## Slide 1" in res.data.extracted_text
        assert "Titel" in res.data.extracted_text
        assert "Hallo" in res.data.extracted_text






