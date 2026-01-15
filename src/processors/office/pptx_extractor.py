from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List

from src.core.exceptions import ProcessingError
from src.core.models.office import OfficeTextContent
from ._common import sanitize_filename, ensure_dir


@dataclass(slots=True)
class PptxExtraction:
    markdown: str
    image_paths: List[str]
    text_contents: List[OfficeTextContent]


def extract_pptx_to_markdown(input_path: Path, images_dir: Path) -> PptxExtraction:
    """Extrahiert PPTX nach Markdown + Embedded-Images (ohne Slide-Rendering).

    Pipeline A ist bewusst „python-only“:
    - Textboxen werden als Text extrahiert
    - Bilder werden als Files extrahiert und im Markdown referenziert
    """
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ProcessingError(f"python-pptx ist nicht installiert: {e}")

    ensure_dir(images_dir)
    pres = Presentation(str(input_path))

    lines: List[str] = []
    image_paths: List[str] = []
    text_contents: List[OfficeTextContent] = []

    for i, slide in enumerate(pres.slides):
        lines.append(f"## Slide {i + 1}")
        slide_text_parts: List[str] = []

        for shape in slide.shapes:
            # Text
            try:
                if getattr(shape, "has_text_frame", False) and shape.text:
                    t = str(shape.text).strip()
                    if t:
                        slide_text_parts.append(t)
                        lines.append(t)
            except Exception:
                pass

            # Picture
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # type: ignore[attr-defined]
                    img = shape.image  # type: ignore[attr-defined]
                    blob: bytes = img.blob  # type: ignore[attr-defined]
                    ext = "." + str(getattr(img, "ext", "bin")).lower()
                    filename = sanitize_filename(f"pptx-img-{i + 1}-{len(image_paths) + 1}") + ext
                    out_path = images_dir / filename
                    with open(out_path, "wb") as f:
                        f.write(blob)
                    rel = f"images/{filename}"
                    image_paths.append(rel)
                    lines.append(f"![]({rel})")
            except Exception:
                pass

        slide_text = "\n".join(slide_text_parts).strip()
        text_contents.append(OfficeTextContent(index=i, text=slide_text, kind="pptx_slide"))
        lines.append("")

    markdown = "\n".join(lines).strip() + "\n"
    return PptxExtraction(markdown=markdown, image_paths=image_paths, text_contents=text_contents)






